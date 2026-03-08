
"""
ppt_generator.py - INTERLOG KPI Dashboard
Genera el PowerPoint mensual con branding INTERLOG.
Compatible con app.py del dashboard.
"""

import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# -------------------------------------------------
# Colores
# -------------------------------------------------
TEAL = RGBColor(63, 168, 156)
NARANJA = RGBColor(212, 105, 10)
ROJO = RGBColor(192, 57, 43)
VERDE = RGBColor(39, 174, 96)
BLANCO = RGBColor(255, 255, 255)


# -------------------------------------------------
# Reglas KPI (misma logica que app.py)
# -------------------------------------------------
def _is_out(item):
    return bool(item.get("desvio")) and str(item.get("parametro", "")).upper() == "INTERLOG"


def _kpi(items):
    total = len(items)
    if total == 0:
        return 0.0, 0, 0, 0

    n_out = sum(1 for i in items if _is_out(i))
    n_in = total - n_out
    pct = (n_in / total) * 100

    return pct, n_in, n_out, total


# -------------------------------------------------
# Utilidades
# -------------------------------------------------
def _empresa(items, nombre):
    return [i for i in items if str(i.get("nombre", "")).upper() == nombre]


# -------------------------------------------------
# Generador principal
# -------------------------------------------------
def generar_ppt(lib_items, ofi_items, cm_pre_items, cm_apr_items, mes="MES"):

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # -------------------------------------------------
    # Slide 1 - Portada
    # -------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "INTERLOG KPI REPORT"
    run.font.size = Pt(44)
    run.font.bold = True

    p = tf.add_paragraph()
    p.text = mes
    p.font.size = Pt(24)

    # -------------------------------------------------
    # Slide 2 - Resumen KPI
    # -------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    y = 1.5
    for titulo, items in [
        ("LIBERACION FASA", _empresa(lib_items, "FASA")),
        ("LIBERACION FSM", _empresa(lib_items, "FSM")),
        ("OFICIALIZACION FASA", _empresa(ofi_items, "FASA")),
        ("OFICIALIZACION FSM", _empresa(ofi_items, "FSM")),
        ("CM PRESENTADOS", cm_pre_items),
    ]:

        pct, n_in, n_out, total = _kpi(items)

        box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(10), Inches(0.8))
        tf = box.text_frame

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{titulo}: {pct:.1f}%  (IN {n_in} / OUT {n_out} / TOTAL {total})"
        run.font.size = Pt(20)

        y += 0.8

    # -------------------------------------------------
    # Exportar a buffer
    # -------------------------------------------------
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return buf
